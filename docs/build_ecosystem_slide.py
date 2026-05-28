"""Build a single-slide PPTX of the Integrated AI Development Ecosystem diagram
in the Accenture brand style.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Accenture brand palette (matches build_pptx.py)
PURPLE = RGBColor(0xA1, 0x00, 0xFF)
DEEP   = RGBColor(0x46, 0x00, 0x73)
LIGHT  = RGBColor(0xB4, 0x55, 0xFF)
SOFT   = RGBColor(0xF2, 0xE6, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY1  = RGBColor(0xF4, 0xF4, 0xF4)
GREY2  = RGBColor(0xBF, 0xBF, 0xBF)
GREY3  = RGBColor(0x5C, 0x5C, 0x5C)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, width_pt=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_text(slide, left, top, width, height, text, *,
             font_size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="Calibri", letter_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))
    return box


def add_rect(slide, left, top, width, height, fill=WHITE, line=None, line_width=1.0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(rect, fill)
    if line is None:
        no_line(rect)
    else:
        set_line(rect, line, line_width)
    rect.shadow.inherit = False
    return rect


def add_node(slide, left, top, width, height, label, *,
             fill=DEEP, fg=WHITE, font_size=14, bold=True, line=None, line_width=1.0):
    rect = add_rect(slide, left, top, width, height, fill=fill,
                    line=line, line_width=line_width)
    add_text(slide, left, top, width, height, label,
             font_size=font_size, bold=bold, color=fg,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return rect


def add_connector(slide, x1, y1, x2, y2, color=PURPLE, width_pt=1.5,
                  kind=MSO_CONNECTOR.STRAIGHT):
    line = slide.shapes.add_connector(kind, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


def add_corner_mark(slide, color=PURPLE):
    add_text(
        slide,
        Inches(12.4), Inches(6.55), Inches(0.7), Inches(0.5),
        ">", font_size=28, bold=True, color=color,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM
    )


def add_brand(slide, color=GREY3):
    add_text(
        slide,
        Inches(9.5), Inches(0.25), Inches(3.7), Inches(0.3),
        "ACCENTURE", font_size=10, color=color,
        align=PP_ALIGN.RIGHT, letter_spacing=2.0
    )


def add_eyebrow(slide, top, text):
    add_text(
        slide,
        Inches(0.6), top, Inches(10), Inches(0.35),
        text.upper(), font_size=11, bold=True, color=PURPLE,
        letter_spacing=2.0
    )


def add_accent_bar(slide, top, color=PURPLE, width_in=0.6):
    return add_rect(
        slide, Inches(0.6), top, Inches(width_in), Inches(0.07), fill=color
    )


def add_title_text(slide, top, text, *, color=BLACK, size=30):
    add_text(
        slide,
        Inches(0.6), top, Inches(12.1), Inches(0.7),
        text, font_size=size, bold=True, color=color
    )


def add_subtitle(slide, top, text, *, color=GREY3, size=13, width_in=11.5):
    add_text(
        slide,
        Inches(0.6), top, Inches(width_in), Inches(0.7),
        text, font_size=size, color=color
    )


def add_pagenum(slide, num, color=GREY3):
    add_text(
        slide,
        Inches(0.4), Inches(6.7), Inches(1.0), Inches(0.3),
        num, font_size=10, color=color,
        align=PP_ALIGN.LEFT, letter_spacing=1.0
    )


def build_ecosystem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=WHITE)

    add_brand(slide)
    add_eyebrow(slide, Inches(0.55), "Integrated AI Development Ecosystem")
    add_title_text(slide, Inches(0.85), "How agents, tools, and IDE come together")
    add_accent_bar(slide, Inches(1.5))
    add_subtitle(
        slide, Inches(1.7),
        "AI coding assistants share a common skill, sub-agent and tool layer, "
        "wired into VS Code alongside source control and enterprise systems."
    )

    # ---- Layout (inches) ----
    # Three columns + right-side VS Code zone
    col_a_x = 0.7   # AI assistants
    col_b_x = 3.4   # Skill / SubAgents / Tools
    top_x   = 6.3   # MCP / Git node (top)
    vscode_x = 8.6  # VS Code box
    side_x = 11.9   # JIRA / Conf / ...

    node_w = 2.0
    node_h = 0.85
    gap_y = 0.35

    # Column A (LLM agents) start y
    a_top = 3.3
    # Column B (capabilities) start y
    b_top = 3.3
    # MCP / Git position
    top_y = 2.55

    # ---- Column A: AI assistants ----
    a_labels = [
        ("MSFT Copilot", DEEP, WHITE),
        ("Gemini CLI", DEEP, WHITE),
        ("Claude Code", DEEP, WHITE),
    ]
    a_rects = []
    for i, (label, fill, fg) in enumerate(a_labels):
        y = a_top + i * (node_h + gap_y)
        r = add_node(
            slide, Inches(col_a_x), Inches(y), Inches(node_w), Inches(node_h),
            label, fill=fill, fg=fg, font_size=13
        )
        a_rects.append((col_a_x, y, node_w, node_h))

    # Column A header label
    add_text(
        slide,
        Inches(col_a_x), Inches(a_top - 0.45), Inches(node_w), Inches(0.32),
        "AI ASSISTANTS", font_size=10, bold=True, color=PURPLE,
        align=PP_ALIGN.CENTER, letter_spacing=2.0
    )

    # ---- Column B: Skill / SubAgents / Tools (shared layer) ----
    b_labels = ["Skills", "Sub-agents", "Tools"]
    b_rects = []
    # Group container (soft purple) behind the three boxes
    group_pad = 0.18
    group_left = col_b_x - group_pad
    group_top = b_top - 0.55
    group_w = node_w + group_pad * 2
    group_h = (node_h + gap_y) * (len(b_labels) - 1) + node_h + 0.75
    add_rect(
        slide, Inches(group_left), Inches(group_top),
        Inches(group_w), Inches(group_h),
        fill=SOFT, line=PURPLE, line_width=1.5
    )
    add_text(
        slide,
        Inches(group_left), Inches(group_top + 0.08),
        Inches(group_w), Inches(0.3),
        "SHARED CAPABILITY LAYER", font_size=10, bold=True, color=DEEP,
        align=PP_ALIGN.CENTER, letter_spacing=2.0
    )
    for i, label in enumerate(b_labels):
        y = b_top + i * (node_h + gap_y)
        add_node(
            slide, Inches(col_b_x), Inches(y), Inches(node_w), Inches(node_h),
            label, fill=PURPLE, fg=WHITE, font_size=13
        )
        b_rects.append((col_b_x, y, node_w, node_h))

    # ---- Top node: MCP / Git ----
    mcp_w = 1.9
    mcp_h = 0.85
    add_node(
        slide, Inches(top_x), Inches(top_y), Inches(mcp_w), Inches(mcp_h),
        "MCP  -  Git", fill=DEEP, fg=WHITE, font_size=13
    )

    # ---- VS Code central box ----
    vs_w = 3.0
    vs_h = 4.4
    vs_top = 2.0
    vs_box = add_rect(
        slide, Inches(vscode_x), Inches(vs_top), Inches(vs_w), Inches(vs_h),
        fill=DEEP, line=PURPLE, line_width=2.0
    )
    add_text(
        slide,
        Inches(vscode_x), Inches(vs_top + 0.4), Inches(vs_w), Inches(0.6),
        "VS Code", font_size=28, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    add_text(
        slide,
        Inches(vscode_x), Inches(vs_top + 1.2), Inches(vs_w), Inches(0.4),
        "DEVELOPER IDE", font_size=10, bold=True, color=LIGHT,
        align=PP_ALIGN.CENTER, letter_spacing=2.0
    )
    # Inner LLMs panel
    inner_pad = 0.25
    inner_left = vscode_x + inner_pad
    inner_w = vs_w - inner_pad * 2
    inner_top = vs_top + vs_h - 1.55
    inner_h = 1.3
    add_rect(
        slide, Inches(inner_left), Inches(inner_top),
        Inches(inner_w), Inches(inner_h),
        fill=LIGHT
    )
    add_text(
        slide,
        Inches(inner_left), Inches(inner_top + 0.15),
        Inches(inner_w), Inches(0.35),
        "LLMs", font_size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )
    add_text(
        slide,
        Inches(inner_left), Inches(inner_top + 0.55),
        Inches(inner_w), Inches(0.7),
        "Gemini  -  Anthropic  -  ...", font_size=12, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP
    )

    # ---- Right-side enterprise systems ----
    side_w = 1.55
    side_h = 0.75
    side_labels = [("JIRA", 2.2), ("Conf.", 3.55), ("...", 4.9)]
    for label, sy in side_labels:
        add_node(
            slide, Inches(side_x), Inches(sy), Inches(side_w), Inches(side_h),
            label, fill=DEEP, fg=WHITE, font_size=13
        )
    add_text(
        slide,
        Inches(side_x), Inches(2.2 - 0.45), Inches(side_w), Inches(0.32),
        "ENTERPRISE", font_size=10, bold=True, color=PURPLE,
        align=PP_ALIGN.CENTER, letter_spacing=2.0
    )

    # ---- Connectors ----
    # AI Assistants -> Shared capability layer (each A to each B)
    for (ax, ay, aw, ah) in a_rects:
        ax_end = ax + aw  # right edge of A
        ay_mid = ay + ah / 2
        for (bx, by, bw, bh) in b_rects:
            by_mid = by + bh / 2
            add_connector(
                slide,
                Inches(ax_end), Inches(ay_mid),
                Inches(bx), Inches(by_mid),
                color=PURPLE, width_pt=1.0
            )

    # Shared capability layer -> VS Code (single bundled connector from group right edge)
    group_right = group_left + group_w
    group_mid_y = group_top + group_h / 2
    add_connector(
        slide,
        Inches(group_right), Inches(group_mid_y),
        Inches(vscode_x), Inches(vs_top + vs_h / 2),
        color=PURPLE, width_pt=2.0
    )

    # MCP/Git -> VS Code (top)
    add_connector(
        slide,
        Inches(top_x + mcp_w), Inches(top_y + mcp_h / 2),
        Inches(vscode_x + vs_w / 2), Inches(vs_top),
        color=PURPLE, width_pt=2.0
    )
    # MCP/Git also wires into shared capability layer (per diagram)
    add_connector(
        slide,
        Inches(top_x), Inches(top_y + mcp_h / 2),
        Inches(group_left + group_w / 2), Inches(group_top),
        color=PURPLE, width_pt=1.0
    )

    # VS Code -> Enterprise systems (right-side)
    for label, sy in side_labels:
        add_connector(
            slide,
            Inches(vscode_x + vs_w), Inches(sy + side_h / 2),
            Inches(side_x), Inches(sy + side_h / 2),
            color=PURPLE, width_pt=1.5
        )

    # Footnote / legend
    add_text(
        slide,
        Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.32),
        "AI assistants share a common skills, sub-agents, and tools layer that "
        "wires into VS Code through MCP / Git, alongside JIRA, Confluence, and "
        "other enterprise systems.",
        font_size=11, color=GREY3
    )

    add_pagenum(slide, "01")
    add_corner_mark(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_ecosystem_slide(prs)
    out = "/Users/rajat.a.ahuja/Dev/FreeRouter/docs/Integrated-AI-Development-Ecosystem.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
